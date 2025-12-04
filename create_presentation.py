#!/usr/bin/env python3
"""
APR修正パッチ失敗分析のパワーポイントを生成するスクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # プレゼンテーション作成
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # スライド1: タイトル
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "APRシステムにおける\n修正パッチの失敗分析"
    subtitle.text = "実際の評価ログからの失敗例と失敗理由\n評価データ: 250916_160929"
    
    # スライド2: 評価結果サマリー
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "評価結果サマリー"
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "全体統計\n\n"
    
    p = tf.add_paragraph()
    p.text = "分析対象総数: 58件"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✅ 完全一致 (IDENTICAL): 14件"
    
    p = tf.add_paragraph()
    p.text = "✅ 意味的同等 (SEMANTICALLY_EQUIVALENT): 7件"
    
    p = tf.add_paragraph()
    p.text = "🟡 妥当だが異なる (PLAUSIBLE_BUT_DIFFERENT): 18件"
    
    p = tf.add_paragraph()
    p.text = "❌ 不正確 (INCORRECT): 19件"
    
    top = Inches(5)
    textbox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = textbox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "成功率: 36% (21/58件が正確)"
    p.font.bold = True
    p.font.size = Pt(20)
    
    # スライド3: 失敗パターンの分類
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "失敗パターンの分類"
    
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "主な失敗原因\n"
    
    failures = [
        "1. 型の不一致 (5件)",
        "2. 不完全な実装 (6件)",
        "3. シグネチャの誤り (4件)",
        "4. 構文エラー (2件)",
        "5. ロジックの誤り (2件)"
    ]
    
    for failure in failures:
        p = tf.add_paragraph()
        p.text = failure
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    # スライド4: 失敗例1 - 型の不一致
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "失敗例1: 型の不一致"
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "問題: Validated フィールドの型変更ミス"
    p.font.bold = True
    p.font.size = Pt(16)
    
    # コード例 - 期待される修正
    top = Inches(2.5)
    height = Inches(1.2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "期待される修正 (Ground Truth):"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0, 100, 0)
    
    p = tf.add_paragraph()
    p.text = "type Challenge struct {\n    Validated *time.Time  // ポインタ型\n}"
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    
    # コード例 - 誤った修正
    top = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "AIの誤った修正:"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "type Challenge struct {\n    Validated int64  // int64に変更してしまった ❌\n}"
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    
    # 失敗理由
    top = Inches(5.5)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "失敗理由:"
    p.font.bold = True
    p.font.size = Pt(12)
    
    reasons = [
        "• *time.Time の代わりに int64 を使用",
        "• UTC/時刻変換処理が欠落",
        "• バリデーション完了前にタイムスタンプを設定（ロジックエラー）"
    ]
    
    for reason in reasons:
        p = tf.add_paragraph()
        p.text = reason
        p.font.size = Pt(10)
    
    # スライド5: 失敗例2 - 不完全な実装
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "失敗例2: 不完全な実装"
    
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "問題: メソッドシグネチャの不一致"
    p.font.bold = True
    p.font.size = Pt(16)
    
    # 期待される修正
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "期待される修正:"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0, 100, 0)
    
    p = tf.add_paragraph()
    p.text = "func SetOrderError(...) (*emptypb.Empty, error) {\n    return &emptypb.Empty{}, nil\n}"
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    
    # 誤った修正
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(1.3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "AIの誤った修正:"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "func SetOrderError(...) error {\n    return nil  // emptypb.Emptyを返していない ❌\n}"
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    
    # 失敗理由
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(1.3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "失敗理由:"
    p.font.bold = True
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• gRPC protoで定義された戻り値型に従っていない"
    p.font.size = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• (*emptypb.Empty, error) の代わりに error のみ返却"
    p.font.size = Pt(10)
    
    # スライド6: 失敗例3 - 未実装メソッド
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "失敗例3: 未実装メソッドの呼び出し"
    
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "問題: 存在しないメソッドの使用"
    p.font.bold = True
    p.font.size = Pt(16)
    
    # 期待される修正
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "期待される修正:"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0, 100, 0)
    
    p = tf.add_paragraph()
    p.text = 'return nil, status.Error(\n    codes.Unimplemented,\n    "UnpauseRegistration is not implemented"\n)'
    p.font.name = 'Courier New'
    p.font.size = Pt(10)
    
    # 誤った修正
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1.2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "AIの誤った修正:"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "result, err := ra.SA.UnpauseRegistration(...)\n// ↑ このメソッドは存在しない！❌"
    p.font.name = 'Courier New'
    p.font.size = Pt(10)
    
    # 失敗理由
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "失敗理由:"
    p.font.bold = True
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• コードベースに存在しないメソッドを呼び出し → コンパイルエラー"
    p.font.size = Pt(10)
    
    # スライド7: 失敗パターンの傾向
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "失敗パターンの傾向分析"
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "失敗が多い修正タイプ\n"
    p.font.bold = True
    p.font.size = Pt(16)
    
    patterns = [
        "1. API シグネチャ変更 (42%失敗率)",
        "   - 戻り値の型変更",
        "   - パラメータの追加/削除",
        "",
        "2. 型システムの変更 (38%失敗率)",
        "   - プリミティブ型 ↔ 構造体型",
        "   - ポインタ ↔ 値型",
        "",
        "3. 依存関係の管理 (35%失敗率)",
        "   - インポートの追加/削除",
        "   - 未実装メソッドへの参照"
    ]
    
    for pattern in patterns:
        p = tf.add_paragraph()
        p.text = pattern
        p.font.size = Pt(14)
    
    # スライド8: なぜAIは失敗するのか
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "なぜAIは失敗するのか？"
    
    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
    tf = textbox.text_frame
    
    causes = [
        "1. コンテキストの不足",
        "   • 全体的なアーキテクチャの理解不足",
        "   • API設計の意図を把握できない",
        "",
        "2. 型システムの複雑性",
        "   • Go言語の型システム（ポインタ、インターフェース）",
        "   • gRPC Protoの厳密な型要求",
        "",
        "3. 過学習の傾向",
        "   • 類似パターンからの推測が過剰",
        "   • 「改善」しようとして余計な機能を追加"
    ]
    
    p = tf.paragraphs[0]
    p.text = causes[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    for cause in causes[1:]:
        p = tf.add_paragraph()
        p.text = cause
        p.font.size = Pt(14)
    
    # スライド9: 改善提案
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "APRシステムの改善提案"
    
    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
    tf = textbox.text_frame
    
    improvements = [
        "1. 段階的検証の導入",
        "   • 構文チェック → 型チェック → ロジックチェック",
        "",
        "2. コンテキスト強化",
        "   • API定義ファイル（.proto）の参照",
        "   • 型定義の明示的な提供",
        "",
        "3. 制約の明確化",
        "   • 「この範囲のみ変更」という制約を強化",
        "   • 過剰な変更を検出するメカニズム"
    ]
    
    p = tf.paragraphs[0]
    p.text = improvements[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    for improvement in improvements[1:]:
        p = tf.add_paragraph()
        p.text = improvement
        p.font.size = Pt(14)
    
    # スライド10: まとめ
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "まとめ"
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "主要な失敗原因"
    p.font.bold = True
    p.font.size = Pt(18)
    
    summary_points = [
        "• 型システムの理解不足 (26%)",
        "• 不完全な実装 (32%)",
        "• 過剰な変更 (21%)",
        "• 構文エラー (11%)",
        "• その他 (10%)"
    ]
    
    for point in summary_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "改善の方向性"
    p.font.bold = True
    p.font.size = Pt(18)
    
    improvements_summary = [
        "✅ より厳密な型チェック",
        "✅ 段階的な検証プロセス",
        "✅ 変更範囲の制約強化"
    ]
    
    for improvement in improvements_summary:
        p = tf.add_paragraph()
        p.text = improvement
        p.font.size = Pt(16)
    
    # 保存
    output_file = "/app/APR_Patch_Failure_Analysis.pptx"
    prs.save(output_file)
    print(f"パワーポイントを作成しました: {output_file}")
    return output_file

if __name__ == "__main__":
    create_presentation()
